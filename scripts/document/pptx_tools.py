#!/usr/bin/env python3
"""
PPTX 文档标准化工具集 v3.1.0

将 4 个独立 PPTX 脚本合并为一个多子命令工具，直接函数调用，不再依赖 subprocess。

子命令:
    font    - 字体统一为微软雅黑（处理所有文本框、表格、母版）
    format  - 文本格式修复（引号、标点、中文单位转标准符号）
    table   - 表格样式设置（标题行、镶边行、首列）
    all     - 一键标准化：依次执行 format -> font -> table

单文件用法（向后兼容）:
    python3 pptx_tools.py <subcommand> [file...]
    python3 pptx_tools.py font presentation.pptx
    python3 pptx_tools.py format file1.pptx file2.pptx
    python3 pptx_tools.py all presentation.pptx

批量并行用法（v3.1+）:
    # --batch FILE : JSONL 每行一个任务（file/subcommand/options）
    python3 pptx_tools.py --batch tasks.jsonl --workers 8
    # JSONL 例子（tasks.jsonl）:
    #   {"file":"/a/x.pptx","subcommand":"font","options":{"do_backup":true}}
    #   {"file":"/a/y.pptx","subcommand":"all","options":{}}
    #   {"file":"/a/z.pptx","subcommand":"all","options":{"phases":"format,font"}}

    # --workers N : ThreadPool 并发度（默认 min(cpu,8)；0=串行）
    python3 pptx_tools.py all *.pptx --workers 4

    # --phases LIST / --defer PHASE : 只对 `all` 子命令有效，控制内部步骤
    python3 pptx_tools.py all x.pptx --phases format,font
    python3 pptx_tools.py all x.pptx --defer table   # 跳过 table 阶段
    # phase 名: format | font | table

    # --fanout-evidence FILE : 写入 PID/线程/任务清单（铁律 #1 真并行 evidence）
    python3 pptx_tools.py all *.pptx --workers 8 \
        --fanout-evidence /tmp/pptx-fanout-evidence.txt

作者: tianli
版本: 3.1.0
日期: 2026-05-23
"""

import argparse
import json
import os
import shutil
import sys
import threading
import time
import traceback
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

# ── lib 路径 ──────────────────────────────────────────────────────────
sys.path.insert(0, str(Path(__file__).parent.parent.parent / "lib"))
sys.path.insert(0, str(Path.home() / "Dev" / "tools" / "dev" / "lib"))  # canonical 5 modules
from finder import get_input_files
from progress import ProgressTracker
from dockit.text import fix_punctuation, fix_quotes, fix_units

# ── 第三方依赖 ────────────────────────────────────────────────────────
try:
    from lxml import etree
    from pptx import Presentation
    from pptx.dml.color import RGBColor  # noqa: F401
    from pptx.oxml.ns import qn
    from pptx.util import Pt  # noqa: F401
except ImportError:
    print("❌ 错误: 缺少 python-pptx 或 lxml 库")
    print("💡 请运行: pip install python-pptx lxml")
    sys.exit(1)


# =====================================================================
#  共用工具函数
# =====================================================================

def show_message(msg_type, message):
    """显示格式化消息"""
    icons = {"success": "✅", "error": "❌", "warning": "⚠️", "info": "ℹ️", "processing": "🔄"}
    icon = icons.get(msg_type, "ℹ️")
    print(f"{icon} {message}")


def backup_file(file_path):
    """备份原始文件"""
    backup_path = f"{file_path}.backup"
    try:
        shutil.copy2(file_path, backup_path)
        show_message("info", f"已备份原文件: {Path(backup_path).name}")
        return backup_path
    except Exception as e:
        show_message("warning", f"备份文件失败: {e}")
        return None


# =====================================================================
#  子命令: font — 字体统一为微软雅黑
# =====================================================================

# 目标字体
TARGET_FONT = "Microsoft YaHei"

# XML 命名空间
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def font_set_for_run(run, font_name=TARGET_FONT):
    """
    为 run 设置字体（强制 XML 级别设置）

    Args:
        run: pptx run 对象
        font_name: 字体名称
    """
    try:
        # 1. 使用 API 设置
        run.font.name = font_name

        # 2. 强制 XML 级别设置 - 直接操作 rPr 元素
        rPr = run._r.get_or_add_rPr()

        # 查找或创建 a:latin 元素（西文字体）
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", font_name)

        # 查找或创建 a:ea 元素（东亚字体 - 中文）
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font_name)

        # 查找或创建 a:cs 元素（复杂脚本字体）
        cs = rPr.find(qn("a:cs"))
        if cs is None:
            cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", font_name)

    except Exception:
        # 某些 run 可能没有字体属性
        pass


def font_set_paragraph_default(paragraph, font_name=TARGET_FONT):
    """
    设置段落的默认字体属性（defRPr）

    Args:
        paragraph: pptx paragraph 对象
        font_name: 字体名称
    """
    try:
        pPr = paragraph._p.get_or_add_pPr()

        # 查找或创建 defRPr（默认文本属性）
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))

        # 设置 latin 字体
        latin = defRPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(defRPr, qn("a:latin"))
        latin.set("typeface", font_name)

        # 设置 ea 字体（东亚）
        ea = defRPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(defRPr, qn("a:ea"))
        ea.set("typeface", font_name)

        # 设置 cs 字体
        cs = defRPr.find(qn("a:cs"))
        if cs is None:
            cs = etree.SubElement(defRPr, qn("a:cs"))
        cs.set("typeface", font_name)

    except Exception:
        pass


def font_set_endParaRPr(paragraph, font_name=TARGET_FONT):
    """
    设置段落结束符的字体属性（endParaRPr）

    Args:
        paragraph: pptx paragraph 对象
        font_name: 字体名称
    """
    try:
        endParaRPr = paragraph._p.find(qn("a:endParaRPr"))
        if endParaRPr is not None:
            # 设置 latin 字体
            latin = endParaRPr.find(qn("a:latin"))
            if latin is None:
                latin = etree.SubElement(endParaRPr, qn("a:latin"))
            latin.set("typeface", font_name)

            # 设置 ea 字体
            ea = endParaRPr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(endParaRPr, qn("a:ea"))
            ea.set("typeface", font_name)

            # 设置 cs 字体
            cs = endParaRPr.find(qn("a:cs"))
            if cs is None:
                cs = etree.SubElement(endParaRPr, qn("a:cs"))
            cs.set("typeface", font_name)
    except Exception:
        pass


def font_process_text_frame(text_frame, stats):
    """
    处理文本框中的所有段落和 run（字体设置）

    Args:
        text_frame: pptx text_frame 对象
        stats: 统计字典
    """
    for paragraph in text_frame.paragraphs:
        # 设置段落默认字体
        font_set_paragraph_default(paragraph, TARGET_FONT)
        # 设置段落结束符字体
        font_set_endParaRPr(paragraph, TARGET_FONT)

        # 处理每个 run
        for run in paragraph.runs:
            font_set_for_run(run, TARGET_FONT)
            stats["font_processed_runs"] += 1


def font_process_table(table, stats):
    """
    处理表格中的所有单元格（字体设置）

    Args:
        table: pptx table 对象
        stats: 统计字典
    """
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                font_process_text_frame(cell.text_frame, stats)
                stats["font_processed_tables"] += 1


def font_process_shape(shape, stats):
    """
    处理单个形状（字体设置）

    Args:
        shape: pptx shape 对象
        stats: 统计字典
    """
    # 处理有文本框的形状
    if shape.has_text_frame:
        font_process_text_frame(shape.text_frame, stats)
        stats["font_processed_shapes"] += 1

    # 处理表格
    if shape.has_table:
        font_process_table(shape.table, stats)

    # 处理组合形状中的子形状
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            font_process_shape(sub_shape, stats)


def font_process_slide(slide, stats):
    """
    处理单个幻灯片（字体设置）

    Args:
        slide: pptx slide 对象
        stats: 统计字典
    """
    for shape in slide.shapes:
        font_process_shape(shape, stats)


def font_process_slide_master(slide_master, stats):
    """
    处理幻灯片母版（字体设置）

    Args:
        slide_master: pptx slide_master 对象
        stats: 统计字典
    """
    # 处理母版中的形状
    for shape in slide_master.shapes:
        font_process_shape(shape, stats)

    # 处理母版的布局
    for layout in slide_master.slide_layouts:
        for shape in layout.shapes:
            font_process_shape(shape, stats)


def font_process_presentation(input_path, do_backup=True):
    """
    格式化 PPT 文档中所有文字的字体为微软雅黑

    Args:
        input_path: 输入文件路径
        do_backup: 是否备份原文件

    Returns:
        bool: 是否成功
    """
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message("error", f"文件不存在: {input_path}")
            return False

        if not input_path.lower().endswith(".pptx"):
            show_message("error", "只支持.pptx格式的文件")
            return False

        show_message("processing", f"正在处理文件: {os.path.basename(input_path)}")

        # 备份原文件
        if do_backup:
            backup_file(input_path)

        # 打开 PPT
        prs = Presentation(input_path)

        # 统计信息
        total_slides = len(prs.slides)
        stats = {
            "font_processed_shapes": 0,
            "font_processed_runs": 0,
            "font_processed_tables": 0,
        }

        show_message("info", f"文档包含 {total_slides} 张幻灯片")

        # 处理幻灯片母版（重要：这里的字体设置会影响整个 PPT）
        show_message("processing", "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            try:
                font_process_slide_master(slide_master, stats)
            except Exception as e:
                show_message("warning", f"处理母版时出错: {e}")

        # 处理所有幻灯片
        show_message("processing", "正在处理幻灯片...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                font_process_slide(slide, stats)
            except Exception as e:
                show_message("warning", f"处理第{i}张幻灯片时出错: {e}")
                continue

        show_message(
            "info",
            f"已处理 {stats['font_processed_shapes']} 个形状, "
            f"{stats['font_processed_runs']} 个文本run",
        )
        if stats["font_processed_tables"] > 0:
            show_message("info", f"已处理 {stats['font_processed_tables']} 个表格单元格")

        # 保存文档
        prs.save(input_path)

        show_message("success", f"字体格式化完成: {os.path.basename(input_path)}")
        show_message("info", f"所有文字已设置为: {TARGET_FONT}")

        return True

    except Exception as e:
        show_message("error", f"处理文件时出错: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: format — 文本格式修复（引号、标点、单位）
# =====================================================================

def format_process_text(text, stats):
    """
    处理文本，应用所有文本转换
    """
    if not text:
        return text

    result, quote_count, _ = fix_quotes(text)
    result, punct_count = fix_punctuation(result)
    result, unit_count = fix_units(result)

    stats["format_quotes"] += quote_count
    stats["format_punctuation"] += punct_count
    stats["format_units"] += unit_count

    return result


def format_process_run(run, stats):
    """
    处理单个 run 的文本（格式修复）
    """
    if run.text:
        original = run.text
        fixed = format_process_text(original, stats)
        if fixed != original:
            run.text = fixed


def format_process_text_frame(text_frame, stats):
    """
    处理文本框中的所有段落和 run（格式修复）
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            format_process_run(run, stats)


def format_process_table(table, stats):
    """
    处理表格中的所有单元格（格式修复）
    """
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                format_process_text_frame(cell.text_frame, stats)


def format_process_shape(shape, stats):
    """
    处理单个形状（格式修复）
    """
    if shape.has_text_frame:
        format_process_text_frame(shape.text_frame, stats)

    if shape.has_table:
        format_process_table(shape.table, stats)

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            format_process_shape(sub_shape, stats)


def format_process_presentation(input_path, do_backup=True):
    """
    处理 PPTX 文件的文本格式（引号、标点、单位）

    Args:
        input_path: 输入文件路径
        do_backup: 是否备份原文件

    Returns:
        bool: 是否成功
    """
    input_p = Path(input_path)

    if not input_p.exists():
        show_message("error", f"文件不存在: {input_path}")
        return False

    if input_p.suffix.lower() != ".pptx":
        show_message("error", "文件必须是.pptx格式")
        return False

    try:
        show_message("processing", f"正在处理文件: {input_p.name}")

        # 备份原文件
        if do_backup:
            backup_file(input_path)

        # 读取文档
        prs = Presentation(input_path)

        # 统计信息
        stats = {"format_quotes": 0, "format_punctuation": 0, "format_units": 0}

        total_slides = len(prs.slides)
        show_message("info", f"文档包含 {total_slides} 张幻灯片")

        # 处理幻灯片母版
        show_message("processing", "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            for shape in slide_master.shapes:
                format_process_shape(shape, stats)
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    format_process_shape(shape, stats)

        # 处理所有幻灯片
        show_message("processing", "正在处理幻灯片...")
        for _i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                format_process_shape(shape, stats)

        # 保存文件（覆盖原文件）
        prs.save(input_path)

        show_message("success", "文本格式修复完成！")
        show_message("info", f"   共替换了 {stats['format_quotes']} 个引号")
        show_message("info", f"   共替换了 {stats['format_punctuation']} 个标点符号")
        show_message("info", f"   共转换了 {stats['format_units']} 个单位")

        return True

    except Exception as e:
        show_message("error", f"处理失败: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: table — 表格样式设置
# =====================================================================

def table_set_style(table):
    """
    设置表格样式选项

    Args:
        table: pptx table 对象

    Returns:
        bool: 是否成功设置
    """
    try:
        # Header Row - 标题行
        table.first_row = True

        # Banded Rows - 镶边行（交替行颜色）
        table.horz_banding = True

        # First Column - 首列
        table.first_col = True

        # 其他可选设置（默认关闭）
        # table.last_row = False      # Total Row - 汇总行
        # table.last_col = False      # Last Column - 末列
        # table.vert_banding = False  # Banded Columns - 镶边列

        return True
    except Exception as e:
        show_message("warning", f"设置表格样式失败: {e}")
        return False


def table_process_shape(shape, stats):
    """
    处理形状，查找表格（表格样式设置）

    Args:
        shape: pptx shape 对象
        stats: 统计字典
    """
    # 处理表格
    if shape.has_table and table_set_style(shape.table):
        stats["table_processed"] += 1

    # 处理组合形状中的子形状
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            table_process_shape(sub_shape, stats)


def table_process_presentation(input_path, do_backup=True):
    """
    处理 PPT 文档中所有表格的样式

    Args:
        input_path: 输入文件路径
        do_backup: 是否备份原文件

    Returns:
        bool: 是否成功
    """
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message("error", f"文件不存在: {input_path}")
            return False

        if not input_path.lower().endswith(".pptx"):
            show_message("error", "只支持.pptx格式的文件")
            return False

        show_message("processing", f"正在处理文件: {os.path.basename(input_path)}")

        # 备份原文件
        if do_backup:
            backup_file(input_path)

        # 打开 PPT
        prs = Presentation(input_path)

        # 统计信息
        total_slides = len(prs.slides)
        stats = {"table_processed": 0}

        show_message("info", f"文档包含 {total_slides} 张幻灯片")

        # 处理所有幻灯片
        show_message("processing", "正在处理表格样式...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                for shape in slide.shapes:
                    table_process_shape(shape, stats)
            except Exception as e:
                show_message("warning", f"处理第{i}张幻灯片时出错: {e}")
                continue

        if stats["table_processed"] > 0:
            show_message("info", f"已处理 {stats['table_processed']} 个表格")
        else:
            show_message("warning", "未找到任何表格")

        # 保存文档
        prs.save(input_path)

        show_message("success", f"表格样式设置完成: {os.path.basename(input_path)}")
        show_message("info", "已启用: Header Row, Banded Rows, First Column")

        return True

    except Exception as e:
        show_message("error", f"处理文件时出错: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: all — 一键标准化（format -> font -> table）
# =====================================================================

# `all` 子命令的内部 phase 切分（用于 --phases / --defer）
ALL_PHASES = ("format", "font", "table")

_PHASE_FUNCS = {
    "format": format_process_presentation,
    "font": font_process_presentation,
    "table": table_process_presentation,
}


def resolve_phases(phases=None, defer=None):
    """
    根据 --phases / --defer 参数解析最终要跑的 phase 列表。

    Args:
        phases: 逗号分隔字符串或 list，None=全部
        defer: 逗号分隔字符串或 list，要跳过的 phase

    Returns:
        list[str]: 有序 phase 名（按 ALL_PHASES 顺序）
    """
    if phases is None:
        chosen = list(ALL_PHASES)
    elif isinstance(phases, str):
        chosen = [p.strip() for p in phases.split(",") if p.strip()]
    else:
        chosen = list(phases)

    if defer:
        if isinstance(defer, str):
            defer_set = {p.strip() for p in defer.split(",") if p.strip()}
        else:
            defer_set = set(defer)
        chosen = [p for p in chosen if p not in defer_set]

    # 校验 + 按 ALL_PHASES 顺序排
    unknown = [p for p in chosen if p not in ALL_PHASES]
    if unknown:
        raise ValueError(f"未知 phase: {unknown}，可选: {ALL_PHASES}")
    return [p for p in ALL_PHASES if p in chosen]


def all_process_presentation(input_path, phases=None, defer=None):
    """
    应用所有 PPTX 标准化处理（直接函数调用，不再使用 subprocess）

    执行顺序：
    1. 文本格式修复（引号、标点、单位）
    2. 字体统一为微软雅黑
    3. 表格样式设置

    Args:
        input_path: 输入文件路径
        phases: 仅跑指定 phase 列表（None=全部）
        defer: 跳过指定 phase 列表

    Returns:
        bool: 是否全部成功
    """
    input_p = Path(input_path)

    # 检查文件是否存在
    if not input_p.exists():
        show_message("error", f"文件不存在: {input_path}")
        return False

    if input_p.suffix.lower() != ".pptx":
        show_message("error", "只支持 .pptx 文件")
        return False

    phase_list = resolve_phases(phases, defer)
    if not phase_list:
        show_message("warning", "phase 列表为空，跳过处理")
        return True

    print("=" * 70)
    print("🚀 开始 PPT 文档标准化处理")
    print("=" * 70)
    print(f"📄 文件: {input_p.name}")
    print(f"🧩 phases: {', '.join(phase_list)}")
    print()

    # 先备份一次（后续步骤不再重复备份）
    backup_file(input_path)

    success_count = 0
    failed_steps = []
    total = len(phase_list)

    for idx, phase in enumerate(phase_list, 1):
        name = f"步骤 {idx}/{total}: {phase}"
        print("\n" + "=" * 70)
        print(f"▶️  {name}")
        print("=" * 70)

        func = _PHASE_FUNCS[phase]
        if func(str(input_p), do_backup=False):
            success_count += 1
            print(f"✅ {name} 完成")
        else:
            failed_steps.append(name)
            print(f"⚠️ {name} 失败（继续执行后续步骤）")

    # 总结
    print("\n" + "=" * 70)
    print("📊 处理总结")
    print("=" * 70)
    print(f"✅ 成功: {success_count}/{total} 个步骤")

    if failed_steps:
        print(f"⚠️ 失败: {len(failed_steps)} 个步骤")
        for step_name in failed_steps:
            print(f"   - {step_name}")
    else:
        print("🎉 所有步骤执行成功！")

    print(f"\n📄 处理完成: {input_p.name}")
    print("=" * 70)

    return len(failed_steps) == 0


# =====================================================================
#  批处理 + 并行（v3.1+）
# =====================================================================

def _dispatch_one(file_path, subcommand, options):
    """
    单任务调度（线程 worker 入口）。

    options 支持的键：
        do_backup: bool（font/format/table 用）
        phases: list 或 csv 字符串（all 用）
        defer: list 或 csv 字符串（all 用）
    """
    options = options or {}
    try:
        if subcommand == "all":
            ok = all_process_presentation(
                file_path,
                phases=options.get("phases"),
                defer=options.get("defer"),
            )
        elif subcommand in ("font", "format", "table"):
            func = {
                "font": font_process_presentation,
                "format": format_process_presentation,
                "table": table_process_presentation,
            }[subcommand]
            ok = func(file_path, do_backup=options.get("do_backup", True))
        else:
            show_message("error", f"未知 subcommand: {subcommand}")
            ok = False
        return {"file": file_path, "subcommand": subcommand, "ok": bool(ok)}
    except Exception as e:
        traceback.print_exc()
        return {"file": file_path, "subcommand": subcommand, "ok": False, "error": str(e)}


def load_batch_jsonl(batch_path):
    """读取 JSONL 任务清单。每行：{"file":"...","subcommand":"...","options":{...}}"""
    tasks = []
    p = Path(batch_path)
    if not p.exists():
        raise FileNotFoundError(f"--batch 文件不存在: {batch_path}")
    with p.open(encoding="utf-8") as f:
        for lineno, raw in enumerate(f, 1):
            line = raw.strip()
            if not line or line.startswith("#"):
                continue
            try:
                obj = json.loads(line)
            except json.JSONDecodeError as e:
                raise ValueError(f"--batch 第 {lineno} 行 JSON 解析失败: {e}") from e
            if "file" not in obj or "subcommand" not in obj:
                raise ValueError(f"--batch 第 {lineno} 行缺 file/subcommand: {line}")
            tasks.append(
                {
                    "file": obj["file"],
                    "subcommand": obj["subcommand"],
                    "options": obj.get("options") or {},
                }
            )
    return tasks


def write_fanout_evidence(path, tasks, workers, start_ts):
    """落地 fan-out evidence（铁律 #1：真并行需 evidence）"""
    try:
        lines = [
            f"# pptx_tools fan-out evidence",
            f"started_at: {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(start_ts))}",
            f"pid: {os.getpid()}",
            f"workers: {workers}",
            f"task_count: {len(tasks)}",
            f"main_thread: {threading.current_thread().name}",
            "",
            "tasks:",
        ]
        for i, t in enumerate(tasks):
            lines.append(f"  [{i}] subcommand={t['subcommand']} file={t['file']}")
        Path(path).write_text("\n".join(lines) + "\n", encoding="utf-8")
        show_message("info", f"fanout-evidence 已写: {path}")
    except Exception as e:
        show_message("warning", f"写 fanout-evidence 失败: {e}")


def run_batch(tasks, workers, fanout_evidence=None):
    """
    并行/串行执行任务清单。

    Args:
        tasks: list[{"file","subcommand","options"}]
        workers: 0=串行；>0=ThreadPool 并发度
        fanout_evidence: 可选 evidence 文件路径

    Returns:
        list[result dict]
    """
    if not tasks:
        show_message("warning", "任务为空")
        return []

    start_ts = time.time()
    if fanout_evidence:
        write_fanout_evidence(fanout_evidence, tasks, workers, start_ts)

    results = []
    if workers == 0 or len(tasks) == 1:
        # 串行
        for t in tasks:
            results.append(_dispatch_one(t["file"], t["subcommand"], t["options"]))
    else:
        # ThreadPool（python-pptx 是 IO+CPU 混合，GIL 下仍能 overlap 多文件 IO/磁盘）
        with ThreadPoolExecutor(max_workers=workers) as ex:
            futs = {
                ex.submit(_dispatch_one, t["file"], t["subcommand"], t["options"]): t
                for t in tasks
            }
            for fut in as_completed(futs):
                results.append(fut.result())

    # 汇总
    elapsed = time.time() - start_ts
    ok_n = sum(1 for r in results if r.get("ok"))
    fail_n = len(results) - ok_n
    print("\n" + "=" * 70)
    print(f"📊 批处理结果: ok={ok_n} fail={fail_n} elapsed={elapsed:.1f}s workers={workers}")
    print("=" * 70)
    for r in results:
        flag = "✅" if r.get("ok") else "❌"
        extra = f" ({r['error']})" if r.get("error") else ""
        print(f"{flag} [{r['subcommand']}] {r['file']}{extra}")
    return results


# =====================================================================
#  CLI 入口
# =====================================================================

def _default_workers():
    """ThreadPool 默认并发度: min(cpu_count, 8)"""
    try:
        n = os.cpu_count() or 4
    except Exception:
        n = 4
    return min(n, 8)


def build_parser():
    """构建 argparse 解析器"""
    parser = argparse.ArgumentParser(
        prog="pptx_tools",
        description="PPTX 文档标准化工具集（v3.1+ 支持 --batch / --workers / --phases / --defer）",
        epilog=(
            "子命令说明:\n"
            "  font    字体统一为微软雅黑\n"
            "  format  文本格式修复（引号、标点、单位）\n"
            "  table   表格样式（标题行、镶边行、首列）\n"
            "  all     一键标准化: format -> font -> table\n"
            "\n"
            "批处理 JSONL 行格式:\n"
            '  {"file":"/a/x.pptx","subcommand":"font","options":{"do_backup":true}}\n'
            '  {"file":"/a/y.pptx","subcommand":"all","options":{"phases":"format,font"}}\n'
            '  {"file":"/a/z.pptx","subcommand":"all","options":{"defer":"table"}}\n'
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "subcommand",
        nargs="?",
        choices=["font", "format", "table", "all"],
        help="子命令: font | format | table | all（使用 --batch 时可省略）",
    )
    parser.add_argument(
        "files",
        nargs="*",
        help="PPTX 文件路径（可多个；不提供则从 Finder 选中获取）",
    )
    # 并行 / 批量 API
    parser.add_argument(
        "--batch",
        metavar="FILE",
        help="JSONL 任务清单（每行 {file, subcommand, options}），与 subcommand/files 互斥",
    )
    parser.add_argument(
        "--workers",
        type=int,
        default=None,
        help=f"ThreadPool 并发度（默认 min(cpu,8)={_default_workers()}；0=串行）",
    )
    parser.add_argument(
        "--phases",
        metavar="LIST",
        help="`all` 子命令的 phase 白名单，逗号分隔（可选: format,font,table）",
    )
    parser.add_argument(
        "--defer",
        metavar="PHASE",
        help="`all` 子命令要跳过的 phase，逗号分隔",
    )
    parser.add_argument(
        "--fanout-evidence",
        metavar="FILE",
        help="写入 fan-out evidence（PID/线程/任务清单），用于铁律 #1 真并行 audit",
    )
    return parser


def main():
    parser = build_parser()

    # 没有参数时显示帮助
    if len(sys.argv) < 2:
        parser.print_help()
        sys.exit(1)

    args = parser.parse_args()
    workers = args.workers if args.workers is not None else _default_workers()

    # ── 模式 1: --batch JSONL ────────────────────────────────────────
    if args.batch:
        if args.subcommand or args.files:
            show_message("warning", "--batch 模式下忽略命令行 subcommand/files")
        try:
            tasks = load_batch_jsonl(args.batch)
        except (FileNotFoundError, ValueError) as e:
            show_message("error", str(e))
            sys.exit(2)
        if not tasks:
            show_message("error", "--batch 文件无有效任务")
            sys.exit(2)
        show_message("info", f"批处理: {len(tasks)} 个任务, workers={workers}")
        results = run_batch(tasks, workers, fanout_evidence=args.fanout_evidence)
        fail_n = sum(1 for r in results if not r.get("ok"))
        sys.exit(0 if fail_n == 0 else 1)

    # ── 模式 2: 传统单/多文件（向后兼容） ────────────────────────────
    if not args.subcommand:
        parser.error("缺少 subcommand（或使用 --batch FILE）")

    files = get_input_files(args.files, expected_ext="pptx")

    if not files:
        show_message("error", "未找到 .pptx 文件")
        print("\n用法: python3 pptx_tools.py <subcommand> [file...]")
        print("  或: python3 pptx_tools.py --batch tasks.jsonl --workers 8")
        print("  或在 Finder 中选择 .pptx 文件后运行")
        sys.exit(1)

    # 把 CLI 参数翻译成 batch tasks → 复用统一调度（顺带启用并行）
    options = {}
    if args.subcommand == "all":
        if args.phases:
            options["phases"] = args.phases
        if args.defer:
            options["defer"] = args.defer
    else:
        options["do_backup"] = True

    tasks = [
        {"file": str(fp), "subcommand": args.subcommand, "options": options}
        for fp in files
    ]

    # 单文件走串行（与旧版输出一致）；多文件按 workers 并行
    effective_workers = 0 if len(tasks) == 1 else workers

    if effective_workers == 0:
        # 保留旧版 ProgressTracker 输出格式
        tracker = ProgressTracker()
        for t in tasks:
            print(f"\n{'=' * 50}")
            print(f"处理文件: {Path(t['file']).name}")
            print("=" * 50)
            r = _dispatch_one(t["file"], t["subcommand"], t["options"])
            if r.get("ok"):
                tracker.add_success()
            else:
                tracker.add_error()
        print(f"\n{'=' * 50}")
        tracker.show_summary("文件处理")
        fail_n = sum(1 for t in tasks if False)  # tracker 已统计；保留兼容
        # 退出码由 tracker 行为主导（保持旧行为：不强制 exit code）
        return

    # 多文件并行
    results = run_batch(tasks, effective_workers, fanout_evidence=args.fanout_evidence)
    fail_n = sum(1 for r in results if not r.get("ok"))
    sys.exit(0 if fail_n == 0 else 1)


if __name__ == "__main__":
    main()
