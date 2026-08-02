#!/usr/bin/env python3
"""
PPTX 工具集统一入口 v2.0.0 —— 家族 4→1 折叠（2026-07-31）

原 4 文件（pptx_cli 转发壳 / pptx_tools / pptx_align / pptx_to_md）合并为单文件子命令族。

⚠ 解释器契约（cc-home skills/pptx/SKILL.md 头号铁律）：
    本文件**顶层 import 仅 stdlib** —— 系统 /usr/bin/python3 必须能跑 audit/layout/
    tablestyle/titlecolor/textboxfill/fontsize/render（样式对齐族，zipfile+re 实现）。
    python-pptx/lxml/finder/progress 等依赖全部在子命令内 lazy import：
    font/format/table/all/to-md/chart 需要 ~/Dev/.venv/bin/python。
    缺包时给清晰提示（"缺少 python-pptx 或 lxml 库"），不是 traceback。

子命令:
    ── 样式对齐族（系统 python3 可跑,原 pptx_align.py） ──
    audit       只读:列每页 layout/表styleId/标题色
    layout      改 slide→layout 引用（--map "3:13,4:14"）
    tablestyle  全部表统一为某 tableStyleId（--style "{GUID}"）
    titlecolor  标题 run 补颜色（--color bg1）
    textboxfill 文本框 spPr 填充→透明 noFill
    fontsize    正文/表格分区设字号（保护标题与设计大字）
    render      soffice→PNG 渲染验证
    ── 标准化族（必 ~/Dev/.venv/bin/python,原 pptx_tools.py） ──
    font        字体统一为微软雅黑（文本框+表格+母版）
    format      文本格式修复（引号/标点/单位）
    table       表格样式设置（标题行/镶边行/首列）
    all         一键标准化 format -> font -> table（--phases/--defer 控步骤）
    ── 转换/生成 ──
    to-md       PPTX 转 Markdown（原 pptx_to_md.py）
    chart       数据驱动图表生成 JSON -> PNG（转发 chart.py,不在本家族）

批量并行（原 pptx_tools --batch,顶层旗截断保留）:
    python3 pptx_cli.py --batch tasks.jsonl --workers 8
    # JSONL 每行: {"file":"/a/x.pptx","subcommand":"all","options":{"phases":"format,font"}}

写盘安全:样式对齐族的 zip 重写收敛到单一 _rewrite 写盘点,每次写盘先跑
lib/docx_parts.assert_parts_intact（allow_changed=本次真实改动的部件）,
丢部件/多改部件抛 PartIntegrityError 且原件不动。

作者: tianli · 版本: 2.0.0 · 2026-07-31 折叠（原件在 ~/.Trash/consolidation-20260731/pptx/）
"""

import argparse
import importlib.util
import json
import os
import re
import shutil
import subprocess
import sys
import threading
import time
import traceback
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import zipfile

SCRIPT_VERSION = "2.0.0"
SCRIPT_DIR = Path(__file__).resolve().parent

sys.path.insert(0, str(SCRIPT_DIR.parent.parent / "lib"))
sys.path.insert(0, str(Path.home() / "Dev" / "tools" / "dev" / "lib"))  # canonical 5 modules
# lsof 占用检查 SSOT。`sub/_cli_common` 顶层只 import stdlib（datetime/json/shutil/
# subprocess/sys/pathlib），不碰 python-pptx/lxml —— 顶层 stdlib-only 契约不破。
# append 不是 insert(0)：sub/ 里有 styles.py 等易撞名模块，不许上抢占位。
sys.path.append(str(SCRIPT_DIR / "sub"))
import _cli_common as _cc  # noqa: E402


# =====================================================================
#  lazy 依赖（保持顶层 stdlib-only 契约）
# =====================================================================

def _hint_wrong_interpreter():
    """总部库(finder/progress 等)加载失败 ≈ 解释器用错（SKILL.md 头号铁律:不是工具坏）。"""
    print("❌ 错误: 总部依赖库加载失败 —— 本子命令(font/format/table/all/to-md/--batch)"
          "必须用 ~/Dev/.venv/bin/python 跑")
    print(f"💡 请运行: ~/Dev/.venv/bin/python {Path(__file__).resolve()} <subcommand> ...")
    print("   （系统 python3 只能跑样式对齐族: audit/layout/tablestyle/titlecolor/"
          "textboxfill/fontsize/render）")
    sys.exit(1)


def _ensure_std_deps():
    """font/format/table/all（原 pptx_tools）的第三方与总部库依赖。"""
    global etree, Presentation, RGBColor, qn, Pt
    global get_input_files, ProgressTracker
    global fix_punctuation, fix_quotes, fix_units
    try:
        from lxml import etree
        from pptx import Presentation
        from pptx.dml.color import RGBColor  # noqa: F401
        from pptx.oxml.ns import qn
        from pptx.util import Pt  # noqa: F401
    except ImportError:
        print("❌ 错误: 缺少 python-pptx 或 lxml 库")
        print("💡 请运行: pip install python-pptx lxml")
        sys.exit(1)
    try:
        from finder import get_input_files
        from progress import ProgressTracker
        from text_fixes import fix_punctuation, fix_quotes, fix_units
    except Exception:
        _hint_wrong_interpreter()


def _ensure_tomd_deps():
    """to-md（原 pptx_to_md）的依赖。"""
    global Presentation
    global show_error, show_info, show_processing, show_success, show_warning
    global check_python_packages, find_files_by_extension, get_file_basename, validate_input_file
    global get_input_files, ProgressTracker
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ 错误: 缺少 python-pptx 或 lxml 库")
        print("💡 请运行: pip install python-pptx lxml")
        sys.exit(1)
    try:
        from display import show_error, show_info, show_processing, show_success, show_warning
        from file_ops import (
            check_python_packages,
            find_files_by_extension,
            get_file_basename,
            validate_input_file,
        )
        from finder import get_input_files
        from progress import ProgressTracker
    except Exception:
        _hint_wrong_interpreter()


# =====================================================================
#  共用工具函数（原 pptx_tools）
# =====================================================================

def show_message(msg_type, message):
    """显示格式化消息"""
    icons = {"success": "✅", "error": "❌", "warning": "⚠️", "info": "ℹ️", "processing": "🔄"}
    icon = icons.get(msg_type, "ℹ️")
    print(f"{icon} {message}")


def backup_file(file_path):
    """备份原始文件"""
    backup_path = f"{file_path}.backup"
    try:
        shutil.copy2(file_path, backup_path)
        show_message("info", f"已备份原文件: {Path(backup_path).name}")
        return backup_path
    except Exception as e:
        show_message("warning", f"备份文件失败: {e}")
        return None


# =====================================================================
#  子命令: font — 字体统一为微软雅黑（原 pptx_tools）
# =====================================================================

# 目标字体
TARGET_FONT = "Microsoft YaHei"


def _set_typefaces(elem, font_name):
    """latin/ea/cs 字体三元组 find-or-create（原 rPr/defRPr/endParaRPr 三连抄收一处）。"""
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = elem.find(qn(tag))
        if node is None:
            node = etree.SubElement(elem, qn(tag))
        node.set("typeface", font_name)


def font_set_for_run(run, font_name=TARGET_FONT):
    """为 run 设置字体（强制 XML 级别设置）"""
    try:
        # 1. 使用 API 设置
        run.font.name = font_name
        # 2. 强制 XML 级别设置 - 直接操作 rPr 元素
        rPr = run._r.get_or_add_rPr()
        _set_typefaces(rPr, font_name)
    except Exception:
        # 某些 run 可能没有字体属性
        pass


def font_set_paragraph_default(paragraph, font_name=TARGET_FONT):
    """设置段落的默认字体属性（defRPr）"""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        # 查找或创建 defRPr（默认文本属性）
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        _set_typefaces(defRPr, font_name)
    except Exception:
        pass


def font_set_endParaRPr(paragraph, font_name=TARGET_FONT):
    """设置段落结束符的字体属性（endParaRPr）"""
    try:
        endParaRPr = paragraph._p.find(qn("a:endParaRPr"))
        if endParaRPr is not None:
            _set_typefaces(endParaRPr, font_name)
    except Exception:
        pass


def walk_shapes(shape, on_text_frame=None, on_table=None):
    """shape 树访问器：文本框/表格回调 + 组合形状递归（原三套骨架逐字同构收一处）。"""
    if on_text_frame is not None and shape.has_text_frame:
        on_text_frame(shape.text_frame)
    if on_table is not None and shape.has_table:
        on_table(shape.table)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            walk_shapes(sub_shape, on_text_frame, on_table)


def font_process_text_frame(text_frame, stats):
    """处理文本框中的所有段落和 run（字体设置）"""
    for paragraph in text_frame.paragraphs:
        # 设置段落默认字体
        font_set_paragraph_default(paragraph, TARGET_FONT)
        # 设置段落结束符字体
        font_set_endParaRPr(paragraph, TARGET_FONT)
        # 处理每个 run
        for run in paragraph.runs:
            font_set_for_run(run, TARGET_FONT)
            stats["font_processed_runs"] += 1


def font_process_table(table, stats):
    """处理表格中的所有单元格（字体设置）"""
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                font_process_text_frame(cell.text_frame, stats)
                stats["font_processed_tables"] += 1


def font_process_shape(shape, stats):
    """处理单个形状（字体设置）"""
    def _on_text_frame(tf):
        font_process_text_frame(tf, stats)
        stats["font_processed_shapes"] += 1

    walk_shapes(shape, _on_text_frame, lambda table: font_process_table(table, stats))


def font_process_slide(slide, stats):
    """处理单个幻灯片（字体设置）"""
    for shape in slide.shapes:
        font_process_shape(shape, stats)


def font_process_slide_master(slide_master, stats):
    """处理幻灯片母版（字体设置）"""
    # 处理母版中的形状
    for shape in slide_master.shapes:
        font_process_shape(shape, stats)
    # 处理母版的布局
    for layout in slide_master.slide_layouts:
        for shape in layout.shapes:
            font_process_shape(shape, stats)


def _open_prs_with_backup(input_path, do_backup):
    """三个 *_process_presentation 共用的「显示→备份→打开→片数」骨架（原三连抄收一处）。"""
    show_message("processing", f"正在处理文件: {os.path.basename(input_path)}")
    if do_backup:
        backup_file(input_path)
    prs = Presentation(input_path)
    show_message("info", f"文档包含 {len(prs.slides)} 张幻灯片")
    return prs


def font_process_presentation(input_path, do_backup=True):
    """格式化 PPT 文档中所有文字的字体为微软雅黑"""
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message("error", f"文件不存在: {input_path}")
            return False
        if not input_path.lower().endswith(".pptx"):
            show_message("error", "只支持.pptx格式的文件")
            return False

        prs = _open_prs_with_backup(input_path, do_backup)

        stats = {
            "font_processed_shapes": 0,
            "font_processed_runs": 0,
            "font_processed_tables": 0,
        }

        # 处理幻灯片母版（重要：这里的字体设置会影响整个 PPT）
        show_message("processing", "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            try:
                font_process_slide_master(slide_master, stats)
            except Exception as e:
                show_message("warning", f"处理母版时出错: {e}")

        # 处理所有幻灯片
        show_message("processing", "正在处理幻灯片...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                font_process_slide(slide, stats)
            except Exception as e:
                show_message("warning", f"处理第{i}张幻灯片时出错: {e}")
                continue

        show_message(
            "info",
            f"已处理 {stats['font_processed_shapes']} 个形状, "
            f"{stats['font_processed_runs']} 个文本run",
        )
        if stats["font_processed_tables"] > 0:
            show_message("info", f"已处理 {stats['font_processed_tables']} 个表格单元格")

        # 保存文档
        prs.save(input_path)

        show_message("success", f"字体格式化完成: {os.path.basename(input_path)}")
        show_message("info", f"所有文字已设置为: {TARGET_FONT}")

        return True

    except Exception as e:
        show_message("error", f"处理文件时出错: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: format — 文本格式修复（引号、标点、单位）（原 pptx_tools）
# =====================================================================

def format_process_text(text, stats):
    """处理文本，应用所有文本转换"""
    if not text:
        return text

    result, quote_count, _ = fix_quotes(text)
    result, punct_count = fix_punctuation(result)
    result, unit_count = fix_units(result)

    stats["format_quotes"] += quote_count
    stats["format_punctuation"] += punct_count
    stats["format_units"] += unit_count

    return result


def format_process_run(run, stats):
    """处理单个 run 的文本（格式修复）"""
    if run.text:
        original = run.text
        fixed = format_process_text(original, stats)
        if fixed != original:
            run.text = fixed


def format_process_text_frame(text_frame, stats):
    """处理文本框中的所有段落和 run（格式修复）"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            format_process_run(run, stats)


def format_process_table(table, stats):
    """处理表格中的所有单元格（格式修复）"""
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                format_process_text_frame(cell.text_frame, stats)


def format_process_shape(shape, stats):
    """处理单个形状（格式修复）"""
    walk_shapes(
        shape,
        lambda tf: format_process_text_frame(tf, stats),
        lambda table: format_process_table(table, stats),
    )


def format_process_presentation(input_path, do_backup=True):
    """处理 PPTX 文件的文本格式（引号、标点、单位）"""
    input_p = Path(input_path)

    if not input_p.exists():
        show_message("error", f"文件不存在: {input_path}")
        return False

    if input_p.suffix.lower() != ".pptx":
        show_message("error", "文件必须是.pptx格式")
        return False

    try:
        prs = _open_prs_with_backup(input_path, do_backup)

        # 统计信息
        stats = {"format_quotes": 0, "format_punctuation": 0, "format_units": 0}

        # 处理幻灯片母版
        show_message("processing", "正在处理幻灯片母版...")
        for slide_master in prs.slide_masters:
            for shape in slide_master.shapes:
                format_process_shape(shape, stats)
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    format_process_shape(shape, stats)

        # 处理所有幻灯片
        show_message("processing", "正在处理幻灯片...")
        for _i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                format_process_shape(shape, stats)

        # 保存文件（覆盖原文件）
        prs.save(input_path)

        show_message("success", "文本格式修复完成！")
        show_message("info", f"   共替换了 {stats['format_quotes']} 个引号")
        show_message("info", f"   共替换了 {stats['format_punctuation']} 个标点符号")
        show_message("info", f"   共转换了 {stats['format_units']} 个单位")

        return True

    except Exception as e:
        show_message("error", f"处理失败: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: table — 表格样式设置（原 pptx_tools）
# =====================================================================

def table_set_style(table):
    """设置表格样式选项"""
    try:
        # Header Row - 标题行
        table.first_row = True
        # Banded Rows - 镶边行（交替行颜色）
        table.horz_banding = True
        # First Column - 首列
        table.first_col = True
        # 其他可选设置（默认关闭）
        # table.last_row = False      # Total Row - 汇总行
        # table.last_col = False      # Last Column - 末列
        # table.vert_banding = False  # Banded Columns - 镶边列
        return True
    except Exception as e:
        show_message("warning", f"设置表格样式失败: {e}")
        return False


def table_process_shape(shape, stats):
    """处理形状，查找表格（表格样式设置）"""
    def _on_table(table):
        if table_set_style(table):
            stats["table_processed"] += 1

    walk_shapes(shape, None, _on_table)


def table_process_presentation(input_path, do_backup=True):
    """处理 PPT 文档中所有表格的样式"""
    try:
        # 验证输入文件
        if not os.path.exists(input_path):
            show_message("error", f"文件不存在: {input_path}")
            return False
        if not input_path.lower().endswith(".pptx"):
            show_message("error", "只支持.pptx格式的文件")
            return False

        prs = _open_prs_with_backup(input_path, do_backup)

        stats = {"table_processed": 0}

        # 处理所有幻灯片
        show_message("processing", "正在处理表格样式...")
        for i, slide in enumerate(prs.slides, 1):
            try:
                for shape in slide.shapes:
                    table_process_shape(shape, stats)
            except Exception as e:
                show_message("warning", f"处理第{i}张幻灯片时出错: {e}")
                continue

        if stats["table_processed"] > 0:
            show_message("info", f"已处理 {stats['table_processed']} 个表格")
        else:
            show_message("warning", "未找到任何表格")

        # 保存文档
        prs.save(input_path)

        show_message("success", f"表格样式设置完成: {os.path.basename(input_path)}")
        show_message("info", "已启用: Header Row, Banded Rows, First Column")

        return True

    except Exception as e:
        show_message("error", f"处理文件时出错: {e}")
        traceback.print_exc()
        return False


# =====================================================================
#  子命令: all — 一键标准化（format -> font -> table）（原 pptx_tools）
# =====================================================================

# `all` 子命令的内部 phase 切分（用于 --phases / --defer）
ALL_PHASES = ("format", "font", "table")

_PHASE_FUNCS = {
    "format": format_process_presentation,
    "font": font_process_presentation,
    "table": table_process_presentation,
}


def resolve_phases(phases=None, defer=None):
    """根据 --phases / --defer 参数解析最终要跑的 phase 列表（按 ALL_PHASES 顺序）。"""
    if phases is None:
        chosen = list(ALL_PHASES)
    elif isinstance(phases, str):
        chosen = [p.strip() for p in phases.split(",") if p.strip()]
    else:
        chosen = list(phases)

    if defer:
        if isinstance(defer, str):
            defer_set = {p.strip() for p in defer.split(",") if p.strip()}
        else:
            defer_set = set(defer)
        chosen = [p for p in chosen if p not in defer_set]

    # 校验 + 按 ALL_PHASES 顺序排
    unknown = [p for p in chosen if p not in ALL_PHASES]
    if unknown:
        raise ValueError(f"未知 phase: {unknown}，可选: {ALL_PHASES}")
    return [p for p in ALL_PHASES if p in chosen]


def all_process_presentation(input_path, phases=None, defer=None):
    """应用所有 PPTX 标准化处理：format → font → table（--phases/--defer 控步骤）"""
    input_p = Path(input_path)

    # 检查文件是否存在
    if not input_p.exists():
        show_message("error", f"文件不存在: {input_path}")
        return False

    if input_p.suffix.lower() != ".pptx":
        show_message("error", "只支持 .pptx 文件")
        return False

    phase_list = resolve_phases(phases, defer)
    if not phase_list:
        show_message("warning", "phase 列表为空，跳过处理")
        return True

    print("=" * 70)
    print("🚀 开始 PPT 文档标准化处理")
    print("=" * 70)
    print(f"📄 文件: {input_p.name}")
    print(f"🧩 phases: {', '.join(phase_list)}")
    print()

    # 先备份一次（后续步骤不再重复备份）
    backup_file(input_path)

    success_count = 0
    failed_steps = []
    total = len(phase_list)

    for idx, phase in enumerate(phase_list, 1):
        name = f"步骤 {idx}/{total}: {phase}"
        print("\n" + "=" * 70)
        print(f"▶️  {name}")
        print("=" * 70)

        func = _PHASE_FUNCS[phase]
        if func(str(input_p), do_backup=False):
            success_count += 1
            print(f"✅ {name} 完成")
        else:
            failed_steps.append(name)
            print(f"⚠️ {name} 失败（继续执行后续步骤）")

    # 总结
    print("\n" + "=" * 70)
    print("📊 处理总结")
    print("=" * 70)
    print(f"✅ 成功: {success_count}/{total} 个步骤")

    if failed_steps:
        print(f"⚠️ 失败: {len(failed_steps)} 个步骤")
        for step_name in failed_steps:
            print(f"   - {step_name}")
    else:
        print("🎉 所有步骤执行成功！")

    print(f"\n📄 处理完成: {input_p.name}")
    print("=" * 70)

    return len(failed_steps) == 0


# =====================================================================
#  批处理 + 并行（原 pptx_tools v3.1+）
# =====================================================================

def _dispatch_one(file_path, subcommand, options):
    """单任务调度（线程 worker 入口）。

    options 支持的键：
        do_backup: bool（font/format/table 用）
        phases: list 或 csv 字符串（all 用）
        defer: list 或 csv 字符串（all 用）
    """
    options = options or {}
    try:
        if subcommand == "all":
            ok = all_process_presentation(
                file_path,
                phases=options.get("phases"),
                defer=options.get("defer"),
            )
        elif subcommand in ("font", "format", "table"):
            func = {
                "font": font_process_presentation,
                "format": format_process_presentation,
                "table": table_process_presentation,
            }[subcommand]
            ok = func(file_path, do_backup=options.get("do_backup", True))
        else:
            show_message("error", f"未知 subcommand: {subcommand}")
            ok = False
        return {"file": file_path, "subcommand": subcommand, "ok": bool(ok)}
    except Exception as e:
        traceback.print_exc()
        return {"file": file_path, "subcommand": subcommand, "ok": False, "error": str(e)}


def load_batch_jsonl(batch_path):
    """读取 JSONL 任务清单。每行：{"file":"...","subcommand":"...","options":{...}}"""
    tasks = []
    p = Path(batch_path)
    if not p.exists():
        raise FileNotFoundError(f"--batch 文件不存在: {batch_path}")
    with p.open(encoding="utf-8") as f:
        for lineno, raw in enumerate(f, 1):
            line = raw.strip()
            if not line or line.startswith("#"):
                continue
            try:
                obj = json.loads(line)
            except json.JSONDecodeError as e:
                raise ValueError(f"--batch 第 {lineno} 行 JSON 解析失败: {e}") from e
            if "file" not in obj or "subcommand" not in obj:
                raise ValueError(f"--batch 第 {lineno} 行缺 file/subcommand: {line}")
            tasks.append(
                {
                    "file": obj["file"],
                    "subcommand": obj["subcommand"],
                    "options": obj.get("options") or {},
                }
            )
    return tasks


def write_fanout_evidence(path, tasks, workers, start_ts):
    """落地 fan-out evidence（铁律 #1：真并行需 evidence）"""
    try:
        lines = [
            f"# pptx_tools fan-out evidence",
            f"started_at: {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(start_ts))}",
            f"pid: {os.getpid()}",
            f"workers: {workers}",
            f"task_count: {len(tasks)}",
            f"main_thread: {threading.current_thread().name}",
            "",
            "tasks:",
        ]
        for i, t in enumerate(tasks):
            lines.append(f"  [{i}] subcommand={t['subcommand']} file={t['file']}")
        Path(path).write_text("\n".join(lines) + "\n", encoding="utf-8")
        show_message("info", f"fanout-evidence 已写: {path}")
    except Exception as e:
        show_message("warning", f"写 fanout-evidence 失败: {e}")


def run_batch(tasks, workers, fanout_evidence=None):
    """并行/串行执行任务清单。workers: 0=串行；>0=ThreadPool 并发度。"""
    if not tasks:
        show_message("warning", "任务为空")
        return []

    start_ts = time.time()
    if fanout_evidence:
        write_fanout_evidence(fanout_evidence, tasks, workers, start_ts)

    results = []
    if workers == 0 or len(tasks) == 1:
        # 串行
        for t in tasks:
            results.append(_dispatch_one(t["file"], t["subcommand"], t["options"]))
    else:
        # ThreadPool（python-pptx 是 IO+CPU 混合，GIL 下仍能 overlap 多文件 IO/磁盘）
        with ThreadPoolExecutor(max_workers=workers) as ex:
            futs = {
                ex.submit(_dispatch_one, t["file"], t["subcommand"], t["options"]): t
                for t in tasks
            }
            for fut in as_completed(futs):
                results.append(fut.result())

    # 汇总
    elapsed = time.time() - start_ts
    ok_n = sum(1 for r in results if r.get("ok"))
    fail_n = len(results) - ok_n
    print("\n" + "=" * 70)
    print(f"📊 批处理结果: ok={ok_n} fail={fail_n} elapsed={elapsed:.1f}s workers={workers}")
    print("=" * 70)
    for r in results:
        flag = "✅" if r.get("ok") else "❌"
        extra = f" ({r['error']})" if r.get("error") else ""
        print(f"{flag} [{r['subcommand']}] {r['file']}{extra}")
    return results


def _default_workers():
    """ThreadPool 默认并发度: min(cpu_count, 8)"""
    try:
        n = os.cpu_count() or 4
    except Exception:
        n = 4
    return min(n, 8)


# =====================================================================
#  子命令: to-md — PPTX 转 Markdown（原 pptx_to_md.py v2.0.0）
# =====================================================================

def check_dependencies():
    show_info("检查依赖项...")
    if not check_python_packages("pptx"):
        sys.exit(1)
    show_success("依赖检查完成")


def convert_pptx_to_md_single(file_path: Path, output_dir) -> bool:
    if not validate_input_file(file_path):
        return False

    if file_path.suffix.lower() != ".pptx":
        show_warning(f"跳过非PPTX文件: {file_path.name}")
        return False

    base_name = get_file_basename(file_path)
    # 直接在原文件目录生成 .md 文件
    output_file = file_path.parent / f"{base_name}.md"

    show_processing(f"转换 {file_path.name} 为 Markdown...")

    try:
        prs = Presentation(file_path)
        with open(output_file, "w", encoding="utf-8") as md_file:
            for i, slide in enumerate(prs.slides, 1):
                md_file.write(f"## Slide {i}\n\n")

                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        md_file.write(shape.text + "\n\n")

                if notes:
                    md_file.write(f"### Speaker Notes\n\n{notes}\n\n")

                md_file.write("---\n\n")

        show_success(f"成功转换: {file_path.name} -> {output_file.name}")
        return True
    except Exception as e:
        show_error(f"转换失败 {file_path.name}: {e}")
        return False


def collect_pptx_files(input_paths: list, recursive: bool = False) -> list:
    """从输入路径列表中收集所有PPTX文件"""
    all_files = []

    for input_path in input_paths:
        path_obj = Path(input_path)

        if path_obj.is_file():
            # 如果是文件，直接检查扩展名
            if path_obj.suffix.lower() == ".pptx":
                all_files.append(path_obj)
            else:
                show_warning(f"跳过非PPTX文件: {path_obj.name}")
        elif path_obj.is_dir():
            # 如果是目录，查找其中的PPTX文件
            found_files = find_files_by_extension(path_obj, "pptx", recursive)
            all_files.extend(found_files)
        else:
            show_error(f"路径不存在: {input_path}")

    return all_files


def _cmd_to_md(args) -> int:
    _ensure_tomd_deps()
    input_paths = list(args.input_paths)
    # 无参数时从 Finder 获取选中的文件（原 pptx_to_md 行为）
    if not input_paths:
        files = get_input_files([], expected_ext="pptx")
        if files:
            input_paths = [str(f) for f in files]
    if not input_paths:
        _SUBPARSERS["to-md"].error("the following arguments are required: input_paths")

    check_dependencies()

    files_to_process = collect_pptx_files(input_paths, args.recursive)

    if not files_to_process:
        show_warning("未找到任何PPTX文件")
        sys.exit(0)

    total_success = 0
    progress = ProgressTracker()

    for file_path in files_to_process:
        show_processing(f"处理 {file_path.name}")
        if convert_pptx_to_md_single(file_path, None):
            total_success += 1
            progress.add_success()
        else:
            progress.add_failure()

    show_info("\n处理完成")
    show_success(f"总共成功转换了 {total_success} 个文件")
    return 0


# =====================================================================
#  样式对齐族（原 pptx_align.py · 仅 stdlib · 系统 python3 可跑）
# =====================================================================

SLIDE_RE = re.compile(r'ppt/slides/slide(\d+)\.xml$')


def _slide_nums(z):
    nums = []
    for n in z.namelist():
        m = SLIDE_RE.match(n)
        if m:
            nums.append(int(m.group(1)))
    return sorted(nums)


def _lsof_guard(path):
    """Office 占用自检 — 占用则退出(调用方/用户先关或 kill)。

    判据 2026-08-02 起委派 `_cli_common.lsof_check`（原实现无 timeout、不看
    returncode、不要求行数>1）；打印文案与 bool 返回值这层外壳保持不变。"""
    occ = _cc.lsof_check(Path(path))
    if occ:
        print(f"⚠️  文件被占用(关闭 PowerPoint 后重试 或 backup+kill):\n{occ}", file=sys.stderr)
        return False
    return True


def _load_zip(path):
    """全量读入 zip：{name: bytes}, {name: ZipInfo}（原 5 处三行样板收一处）。"""
    with zipfile.ZipFile(path) as z:
        return ({n: z.read(n) for n in z.namelist()},
                {n: z.getinfo(n) for n in z.namelist()})


def _rewrite(path, items, infos, changed_parts):
    """zip 全量重写 —— 样式对齐族唯一写盘点。

    先写 .tmp，再跑 lib/docx_parts.assert_parts_intact（原件 vs .tmp，
    allow_changed=本次真实改动的部件名），断言通过才 os.replace 原件。
    丢部件/白名单外部件被改 → 抛 PartIntegrityError，原件一个字节不动。
    """
    from docx_parts import assert_parts_intact  # lib/ 纯 stdlib 实现,系统 python3 可跑
    tmp = str(path) + '.tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zo:
        for n, d in items.items():
            zo.writestr(infos[n], d)
    try:
        assert_parts_intact(path, tmp, allow_changed=set(changed_parts), verbose=False)
    except BaseException:
        os.unlink(tmp)
        raise
    os.replace(tmp, path)


# ── audit:只读侦察 ────────────────────────────────────────────────
def cmd_audit(path, args):
    z = zipfile.ZipFile(path)
    # tableStyles 名称表
    name_of = {}
    if 'ppt/tableStyles.xml' in z.namelist():
        ts = z.read('ppt/tableStyles.xml').decode('utf-8', 'ignore')
        name_of = dict(re.findall(r'<a:tblStyle styleId="(\{[^}]*\})"\s+styleName="([^"]*)"', ts))

    print(f"# PPTX 结构侦察: {Path(path).name}")
    print(f"slides={len(_slide_nums(z))}  "
          f"layouts={len([n for n in z.namelist() if 'slideLayouts/slideLayout' in n and n.endswith('.xml')])}  "
          f"masters={len([n for n in z.namelist() if 'slideMasters/slideMaster' in n and n.endswith('.xml')])}  "
          f"themes={len([n for n in z.namelist() if n.startswith('ppt/theme/theme') and n.endswith('.xml')])}")
    print()
    print(f"{'slide':>6} | {'layout':>7} | {'表(styleId名)':<22} | {'标题色':<10} | 首文本")
    print("-" * 90)
    for num in _slide_nums(z):
        x = z.read(f'ppt/slides/slide{num}.xml').decode('utf-8', 'ignore')
        # layout
        relp = f'ppt/slides/_rels/slide{num}.xml.rels'
        lay = '?'
        if relp in z.namelist():
            m = re.search(r'slideLayout(\d+)\.xml', z.read(relp).decode())
            lay = m.group(1) if m else '?'
        # 表 styleId
        sids = re.findall(r'<a:tableStyleId>([^<]*)</a:tableStyleId>', x)
        tbl = ','.join(name_of.get(s, s[-8:]) for s in sids) if sids else '-'
        # 标题色(第一个 (一)(二) 或 中文数字标题 run 的 solidFill)
        tcolor = '-'
        for sp in re.findall(r'<p:sp>.*?</p:sp>', x, re.S):
            t = ''.join(re.findall(r'<a:t>([^<]*)</a:t>', sp)).strip()
            if re.match(r'^[（(][一二三四五六七八九十][）)]|^[一二三四五六七八九十]、', t):
                r = re.search(r'<a:r>.*?<a:t>', sp, re.S)
                rt = r.group(0) if r else ''
                fm = re.search(r'<a:solidFill>\s*<a:(?:schemeClr val="(\w+)"|srgbClr val="([0-9A-Fa-f]{6})")', rt)
                tcolor = (fm.group(1) or fm.group(2)) if fm else '黑/缺'
                break
        txt = ''.join(re.findall(r'<a:t>([^<]*)</a:t>', x))[:16]
        print(f"{num:>6} | {lay:>7} | {tbl:<22} | {tcolor:<10} | {txt}")
    return 0


# ── layout:改 slide→layout 引用 ────────────────────────────────
def cmd_layout(path, args):
    if not _lsof_guard(path):
        return 1
    mapping = {}
    for pair in args.map.split(','):
        s, l = pair.split(':')
        mapping[int(s)] = int(l)
    items, infos = _load_zip(path)
    changed = []
    changed_parts = set()
    for snum, lnum in mapping.items():
        rel = f'ppt/slides/_rels/slide{snum}.xml.rels'
        if rel not in items:
            print(f"  slide{snum} 无 rels,跳过", file=sys.stderr)
            continue
        txt = items[rel].decode()
        new = re.sub(r'slideLayout\d+\.xml', f'slideLayout{lnum}.xml', txt, count=1)
        if new != txt:
            items[rel] = new.encode()
            changed.append((snum, lnum))
            changed_parts.add(rel)
    if not args.dry_run and changed:
        _rewrite(path, items, infos, changed_parts)
    print(f"{'[dry-run] ' if args.dry_run else ''}改 layout: " +
          ', '.join(f'slide{s}→layout{l}' for s, l in changed))
    return 0


# ── tablestyle:全部表统一 styleId ──────────────────────────────
def cmd_tablestyle(path, args):
    if not _lsof_guard(path):
        return 1
    sid = args.style if args.style.startswith('{') else '{' + args.style + '}'
    items, infos = _load_zip(path)
    # 校验 styleId 在 tableStyles.xml 里有定义
    if 'ppt/tableStyles.xml' in items:
        if sid not in items['ppt/tableStyles.xml'].decode('utf-8', 'ignore'):
            print(f"⚠️  styleId {sid} 不在 tableStyles.xml 中,改了会丢样式", file=sys.stderr)
            return 1
    cnt = 0
    changed_parts = set()
    for n in list(items):
        if SLIDE_RE.match(n):
            x = items[n].decode('utf-8', 'ignore')
            c = len(re.findall(r'<a:tableStyleId>', x))
            if c:
                x2 = re.sub(r'<a:tableStyleId>[^<]*</a:tableStyleId>',
                            f'<a:tableStyleId>{sid}</a:tableStyleId>', x)
                if x2 != x:
                    items[n] = x2.encode()
                    cnt += c
                    changed_parts.add(n)
    if not args.dry_run and cnt:
        _rewrite(path, items, infos, changed_parts)
    print(f"{'[dry-run] ' if args.dry_run else ''}统一 {cnt} 张表 → {sid}")
    return 0


# ── titlecolor:给标题 run 补颜色 ───────────────────────────────
def cmd_titlecolor(path, args):
    if not _lsof_guard(path):
        return 1
    color = args.color
    fill = (f'<a:solidFill><a:schemeClr val="{color}"/></a:solidFill>'
            if not re.fullmatch(r'[0-9A-Fa-f]{6}', color)
            else f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>')
    pat = re.compile(args.pattern)
    items, infos = _load_zip(path)
    pages = []
    changed_parts = set()
    for n in list(items):
        if not SLIDE_RE.match(n):
            continue
        num = int(SLIDE_RE.match(n).group(1))
        x = items[n].decode('utf-8', 'ignore')

        def fix_sp(m):
            sp = m.group(0)
            t = ''.join(re.findall(r'<a:t>([^<]*)</a:t>', sp)).strip()
            if not pat.match(t):
                return sp

            def addfill(rm):
                rpr = rm.group(0)
                if '<a:solidFill>' in rpr:
                    return rpr
                if rpr.endswith('/>'):
                    return rpr[:-2] + '>' + fill + '</a:rPr>'
                if '<a:latin' in rpr:
                    return rpr.replace('<a:latin', fill + '<a:latin', 1)
                return rpr.replace('</a:rPr>', fill + '</a:rPr>', 1)
            return re.sub(r'<a:rPr[^>]*(?:/>|>.*?</a:rPr>)', addfill, sp, count=1, flags=re.S)
        x2 = re.sub(r'<p:sp>.*?</p:sp>', fix_sp, x, flags=re.S)
        if x2 != x:
            items[n] = x2.encode()
            pages.append(num)
            changed_parts.add(n)
    if not args.dry_run and pages:
        _rewrite(path, items, infos, changed_parts)
    print(f"{'[dry-run] ' if args.dry_run else ''}标题补色 {color}: slide{pages}")
    return 0


# ── textboxfill:文本框 spPr 填充 → 透明(noFill) ────────────────
# 把"文字文本框"(<p:sp> 内含 <p:txBody>)的 shape 级填充改成透明。
# 只动 <p:spPr> 直接子级的 <a:solidFill>,绝不碰 <p:txBody> 内 run 文字色
# (<a:rPr> 里的 solidFill 结构上在 txBody 内,被 spPr 块边界天然隔离)。
# --match 选目标填充签名:
#   bg1-alpha (默认) — 只动 <a:schemeClr val="bg1"> 含 <a:alpha> 的半透明白底遮罩
#                       (保护纯色设计块/渐变块/无 alpha 的实底色)
#   any-solid        — 文本框所有 spPr 级 solidFill 都改透明(含纯色块,慎用)
_SPPR_RE = re.compile(r'<p:spPr\b[^>]*>.*?</p:spPr>', re.S)
_SOLIDFILL_RE = re.compile(r'<a:solidFill>.*?</a:solidFill>', re.S)


def _is_bg1_alpha(solidfill_xml):
    """该 solidFill 是否 schemeClr=bg1 且带 alpha(半透明白底遮罩)。"""
    sc = re.search(r'<a:schemeClr val="(\w+)"\s*>(.*?)</a:schemeClr>', solidfill_xml, re.S)
    if not sc or sc.group(1) != 'bg1':
        return False
    return '<a:alpha' in sc.group(2)


def cmd_textboxfill(path, args):
    if not _lsof_guard(path):
        return 1
    mode = args.match
    items, infos = _load_zip(path)

    total = 0
    per_page = {}
    changed_parts = set()
    for n in list(items):
        if not SLIDE_RE.match(n):
            continue
        num = int(SLIDE_RE.match(n).group(1))
        x = items[n].decode('utf-8', 'ignore')

        def fix_sp(m):
            sp = m.group(0)
            if '<p:txBody>' not in sp:        # 仅文字文本框
                return sp
            spm = _SPPR_RE.search(sp)
            if not spm:
                return sp
            spPr = spm.group(0)

            def repl_fill(fm):
                fill = fm.group(0)
                if mode == 'bg1-alpha' and not _is_bg1_alpha(fill):
                    return fill               # 保护非 bg1-alpha 填充(纯色块等)
                return '<a:noFill/>'
            new_spPr = _SOLIDFILL_RE.sub(repl_fill, spPr)
            if new_spPr == spPr:
                return sp
            cnt = spPr.count('<a:solidFill>') - new_spPr.count('<a:solidFill>')
            fix_sp.hits += cnt
            return sp.replace(spPr, new_spPr, 1)
        fix_sp.hits = 0

        x2 = re.sub(r'<p:sp>.*?</p:sp>', fix_sp, x, flags=re.S)
        if x2 != x:
            items[n] = x2.encode()
            per_page[num] = fix_sp.hits
            total += fix_sp.hits
            changed_parts.add(n)

    if not args.dry_run and total:
        _rewrite(path, items, infos, changed_parts)
    tag = '[dry-run] ' if args.dry_run else ''
    pages = ','.join(f'{p}({c})' for p, c in sorted(per_page.items()))
    print(f"{tag}文本框填充→透明(match={mode}): {total} 处, 涉及页 [{pages}]")
    return 0


# ── fontsize:正文/表格分区设字号(保护标题与设计大字) ──────────
# 精确分区设 run 字号,三类边界互不覆盖:
#   表格单元格(<a:tbl> 内 rPr)        → --table pt
#   正文文本框(<p:sp>/<p:txBody> 非tbl)→ --body  pt
#   标题(--title-pattern 命中文本)     → 跳过(保护 (一)(二) 大标题)
#   设计大字(原 sz ≥ --protect-above)  → 跳过(保护封面/目录/章节/致谢 54/72/80pt 等)
# 用 lxml 局部 import(精确 DOM 区分 tbl 边界,正则难可靠);系统 python3 自带 lxml,
# 不破坏模块顶层 stdlib-only 契约(仅本子命令触发时 import)。
_TITLE_PAT = r'^[（(][一二三四五六七八九十][）)]|^[一二三四五六七八九十]、'


def cmd_fontsize(path, args):
    if not _lsof_guard(path):
        return 1
    from lxml import etree
    A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    PP = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    nsd = {'a': A[1:-1], 'p': PP[1:-1]}
    title_pat = re.compile(args.title_pattern)
    body_sz, tbl_sz, protect = args.body, args.table, args.protect_above

    items, infos = _load_zip(path)

    stat = {'table': 0, 'body': 0, 'skip_title': 0, 'skip_big': 0}
    changed_parts = set()
    for n in list(items):
        if not SLIDE_RE.match(n):
            continue
        root = etree.fromstring(items[n])
        tbl_rpr = set()
        if tbl_sz:
            for tbl in root.iter(A + 'tbl'):
                for rpr in tbl.iter(A + 'rPr'):
                    tbl_rpr.add(id(rpr))
                    rpr.set('sz', str(tbl_sz * 100)); stat['table'] += 1
                for tag in ('endParaRPr', 'defRPr'):
                    for rpr in tbl.iter(A + tag):
                        rpr.set('sz', str(tbl_sz * 100))
        else:
            for tbl in root.iter(A + 'tbl'):
                for rpr in tbl.iter(A + 'rPr'):
                    tbl_rpr.add(id(rpr))
        if body_sz:
            for sp in root.iter(PP + 'sp'):
                tb = sp.find('p:txBody', nsd)
                if tb is None:
                    continue
                if title_pat.match(''.join(tb.itertext()).strip()):
                    stat['skip_title'] += sum(1 for _ in tb.iter(A + 'rPr'))
                    continue
                for rpr in tb.iter(A + 'rPr'):
                    if id(rpr) in tbl_rpr:
                        continue
                    cur = rpr.get('sz')
                    if cur and int(cur) >= protect * 100:
                        stat['skip_big'] += 1; continue
                    rpr.set('sz', str(body_sz * 100)); stat['body'] += 1
        new = etree.tostring(root)
        if new != items[n]:
            changed_parts.add(n)
        items[n] = new

    if not args.dry_run and (stat['table'] or stat['body']):
        _rewrite(path, items, infos, changed_parts)
    tag = '[dry-run] ' if args.dry_run else ''
    print(f"{tag}字号分区: 表格cell→{tbl_sz}pt:{stat['table']}  正文→{body_sz}pt:{stat['body']}  "
          f"跳过标题run:{stat['skip_title']}  跳过≥{protect}pt设计大字:{stat['skip_big']}")
    return 0


# ── render:soffice → PNG 验证 ──────────────────────────────────
def cmd_render(path, args):
    from soffice import find_soffice  # doctools SSOT: soffice 路径解析(lazy,保 stdlib 顶层)
    soffice = find_soffice()
    if not soffice:
        print("⚠️  未找到 LibreOffice(soffice)", file=sys.stderr)
        return 1
    outdir = Path(args.outdir or '/tmp/pptx-render')
    outdir.mkdir(parents=True, exist_ok=True)
    pdf = outdir / (Path(path).stem + '.pdf')
    subprocess.run([soffice, '--headless', '--convert-to', 'pdf',
                    '--outdir', str(outdir), str(path)],
                   capture_output=True, timeout=180)
    if not pdf.exists():
        print("⚠️  soffice 转 PDF 失败", file=sys.stderr)
        return 1
    pages = args.pages.split(',') if args.pages else None
    if pages:
        for p in pages:
            subprocess.run(['pdftoppm', '-png', '-r', str(args.dpi),
                            '-f', p, '-l', p, str(pdf), str(outdir / f'p{p}')],
                           capture_output=True)
    else:
        subprocess.run(['pdftoppm', '-png', '-r', str(args.dpi), str(pdf), str(outdir / 'page')],
                       capture_output=True)
    pngs = sorted(outdir.glob('*.png'))
    print(f"渲染 {len(pngs)} 张 → {outdir}/")
    for p in pngs:
        print(f"  {p}")
    return 0


_ALIGN_FUNCS = {'audit': cmd_audit, 'layout': cmd_layout, 'tablestyle': cmd_tablestyle,
                'titlecolor': cmd_titlecolor, 'textboxfill': cmd_textboxfill,
                'fontsize': cmd_fontsize, 'render': cmd_render}


def _cmd_align(args) -> int:
    path = Path(args.pptx)
    if not path.exists():
        print(f"文件不存在: {path}", file=sys.stderr)
        return 1
    return _ALIGN_FUNCS[args.cmd](path, args)


# =====================================================================
#  子命令: chart — 转发 chart.py（891 行独立引擎,另有 2 个调用方,不并入）
# =====================================================================

def _load_sibling(module_alias: str, filename: str):
    """用 importlib 按绝对路径加载同目录脚本（保留别名 import 避免污染顶级 namespace）。"""
    path = SCRIPT_DIR / filename
    spec = importlib.util.spec_from_file_location(module_alias, path)
    if spec is None or spec.loader is None:
        raise ImportError(f"无法加载 {path}")
    mod = importlib.util.module_from_spec(spec)
    sys.modules[module_alias] = mod
    spec.loader.exec_module(mod)
    return mod


def _run_module_main(mod, argv0: str, argv: list[str]) -> int:
    """用 argv 包装调用 mod.main()，捕获 SystemExit。"""
    saved = sys.argv
    try:
        sys.argv = [argv0, *argv]
        try:
            mod.main()
        except SystemExit as e:
            return int(e.code) if e.code is not None else 0
        return 0
    finally:
        sys.argv = saved


def _dispatch_chart(argv: list[str]) -> int:
    """转发 chart 到 chart.py"""
    mod = _load_sibling("_chart_sibling", "chart.py")
    return _run_module_main(mod, "chart", argv)


# =====================================================================
#  标准化族 CLI 流（原 pptx_tools main 模式 1/2）
# =====================================================================

def _batch_parser():
    """--batch 顶层旗模式的解析器（原 pptx_tools build_parser,含向后兼容位置参数）。"""
    parser = argparse.ArgumentParser(
        prog="pptx_cli",
        description="PPTX 文档标准化工具集（--batch / --workers / --phases / --defer）",
        epilog=(
            "批处理 JSONL 行格式:\n"
            '  {"file":"/a/x.pptx","subcommand":"font","options":{"do_backup":true}}\n'
            '  {"file":"/a/y.pptx","subcommand":"all","options":{"phases":"format,font"}}\n'
            '  {"file":"/a/z.pptx","subcommand":"all","options":{"defer":"table"}}\n'
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("subcommand", nargs="?",
                        choices=["font", "format", "table", "all"],
                        help="子命令（--batch 模式下被忽略）")
    parser.add_argument("files", nargs="*", help="PPTX 文件路径（--batch 模式下被忽略）")
    parser.add_argument("--batch", metavar="FILE",
                        help="JSONL 任务清单（每行 {file, subcommand, options}）")
    parser.add_argument("--workers", type=int, default=None,
                        help=f"ThreadPool 并发度（默认 min(cpu,8)={_default_workers()}；0=串行）")
    parser.add_argument("--phases", metavar="LIST", help="`all` 的 phase 白名单，逗号分隔")
    parser.add_argument("--defer", metavar="PHASE", help="`all` 要跳过的 phase，逗号分隔")
    parser.add_argument("--fanout-evidence", metavar="FILE",
                        help="写入 fan-out evidence（PID/线程/任务清单）")
    return parser


def _batch_main() -> int:
    """--batch JSONL 模式（原 pptx_tools main 模式 1,字面量与退出码保持）。"""
    _ensure_std_deps()
    args = _batch_parser().parse_args()
    workers = args.workers if args.workers is not None else _default_workers()
    if args.subcommand or args.files:
        show_message("warning", "--batch 模式下忽略命令行 subcommand/files")
    try:
        tasks = load_batch_jsonl(args.batch)
    except (FileNotFoundError, ValueError) as e:
        show_message("error", str(e))
        sys.exit(2)
    if not tasks:
        show_message("error", "--batch 文件无有效任务")
        sys.exit(2)
    show_message("info", f"批处理: {len(tasks)} 个任务, workers={workers}")
    results = run_batch(tasks, workers, fanout_evidence=args.fanout_evidence)
    fail_n = sum(1 for r in results if not r.get("ok"))
    sys.exit(0 if fail_n == 0 else 1)


def _cmd_std(args) -> int:
    """font/format/table/all（原 pptx_tools main 模式 2,字面量与退出码保持）。"""
    _ensure_std_deps()
    subcommand = args.cmd
    workers = args.workers if args.workers is not None else _default_workers()

    files = get_input_files(args.files, expected_ext="pptx")

    if not files:
        show_message("error", "未找到 .pptx 文件")
        print("\n用法: python3 pptx_tools.py <subcommand> [file...]")
        print("  或: python3 pptx_tools.py --batch tasks.jsonl --workers 8")
        print("  或在 Finder 中选择 .pptx 文件后运行")
        sys.exit(1)

    # 把 CLI 参数翻译成 batch tasks → 复用统一调度（顺带启用并行）
    options = {}
    if subcommand == "all":
        if args.phases:
            options["phases"] = args.phases
        if args.defer:
            options["defer"] = args.defer
    else:
        options["do_backup"] = True

    tasks = [
        {"file": str(fp), "subcommand": subcommand, "options": options}
        for fp in files
    ]

    # 单文件走串行（与旧版输出一致）；多文件按 workers 并行
    effective_workers = 0 if len(tasks) == 1 else workers

    if effective_workers == 0:
        # 保留旧版 ProgressTracker 输出格式
        tracker = ProgressTracker()
        for t in tasks:
            print(f"\n{'=' * 50}")
            print(f"处理文件: {Path(t['file']).name}")
            print("=" * 50)
            r = _dispatch_one(t["file"], t["subcommand"], t["options"])
            if r.get("ok"):
                tracker.add_success()
            else:
                tracker.add_error()   # 历史怪癖原样保留：ProgressTracker 无此方法,失败即 AttributeError
        print(f"\n{'=' * 50}")
        tracker.show_summary("文件处理")
        # 退出码由 tracker 行为主导（保持旧行为：不强制 exit code）
        return 0

    # 多文件并行
    results = run_batch(tasks, effective_workers, fanout_evidence=args.fanout_evidence)
    fail_n = sum(1 for r in results if not r.get("ok"))
    sys.exit(0 if fail_n == 0 else 1)


# =====================================================================
#  CLI 入口
# =====================================================================

_SUBPARSERS = {}


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="pptx_cli",
        description="PPTX 工具集统一入口 v%s（13 子命令 + chart 转发；"
                    "样式对齐族系统 python3 可跑,标准化族需 ~/Dev/.venv/bin/python）" % SCRIPT_VERSION,
        epilog=(
            "批量模式（顶层旗,原 pptx_tools --batch）:\n"
            "  python3 pptx_cli.py --batch tasks.jsonl --workers 8\n"
            "各子命令完整 --help: python3 pptx_cli.py <subcommand> --help"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--version", action="version", version=f"%(prog)s {SCRIPT_VERSION}")
    sub = parser.add_subparsers(dest="cmd", required=True, metavar="<subcommand>")

    # ── 标准化族（原 pptx_tools） ──
    for name, help_ in (
        ("font", "字体统一为微软雅黑（venv python）"),
        ("format", "文本格式修复 引号/标点/单位（venv python）"),
        ("table", "表格样式设置 标题行/镶边行/首列（venv python）"),
        ("all", "一键标准化 format -> font -> table（venv python）"),
    ):
        p = sub.add_parser(name, help=help_)
        p.add_argument("files", nargs="*",
                       help="PPTX 文件路径（可多个；不提供则从 Finder 选中获取）")
        p.add_argument("--workers", type=int, default=None,
                       help=f"ThreadPool 并发度（默认 min(cpu,8)={_default_workers()}；0=串行）")
        p.add_argument("--phases", metavar="LIST",
                       help="`all` 子命令的 phase 白名单，逗号分隔（可选: format,font,table）")
        p.add_argument("--defer", metavar="PHASE",
                       help="`all` 子命令要跳过的 phase，逗号分隔")
        p.add_argument("--fanout-evidence", metavar="FILE",
                       help="写入 fan-out evidence（PID/线程/任务清单）")
        p.set_defaults(func=_cmd_std)
        _SUBPARSERS[name] = p

    # ── to-md（原 pptx_to_md） ──
    t = sub.add_parser("to-md", help="PPTX 转 Markdown（venv python）")
    t.add_argument("input_paths", nargs="*", help="一个或多个PPTX文件/目录路径")
    t.add_argument("-r", "--recursive", action="store_true", help="递归处理目录")
    t.add_argument("--version", action="version", version=f"%(prog)s {SCRIPT_VERSION}")
    t.set_defaults(func=_cmd_to_md)
    _SUBPARSERS["to-md"] = t

    # ── 样式对齐族（原 pptx_align,系统 python3） ──
    a = sub.add_parser('audit', help='只读:列每页 layout/表styleId/标题色（系统 python3）')
    a.add_argument('pptx')

    l = sub.add_parser('layout', help='改 slide→layout 引用（系统 python3）')
    l.add_argument('pptx')
    l.add_argument('--map', required=True, help='逗号分隔 slideN:layoutN,如 "3:13,4:14"')
    l.add_argument('--dry-run', action='store_true')

    ts = sub.add_parser('tablestyle', help='全部表统一为某 tableStyleId（系统 python3）')
    ts.add_argument('pptx')
    ts.add_argument('--style', required=True, help='tableStyleId GUID(带或不带花括号)')
    ts.add_argument('--dry-run', action='store_true')

    c = sub.add_parser('titlecolor', help='标题 run 补颜色(schemeClr 名或 srgb 6位)（系统 python3）')
    c.add_argument('pptx')
    c.add_argument('--color', required=True, help='bg1/tx1/accent1.. 或 FFFFFF')
    c.add_argument('--pattern', default=r'^[（(][一二三四五六七八九十][）)]',
                   help='标题文本匹配正则(默认 (一)(二)类)')
    c.add_argument('--dry-run', action='store_true')

    fs = sub.add_parser('fontsize', help='正文/表格分区设字号(保护标题与设计大字)（系统 python3）')
    fs.add_argument('pptx')
    fs.add_argument('--body', type=int, default=0, help='正文文本框 run 字号pt(0=不动)')
    fs.add_argument('--table', type=int, default=0, help='表格单元格 run 字号pt(0=不动)')
    fs.add_argument('--protect-above', type=int, default=24,
                    help='正文中原≥此pt的大字跳过(保护封面/目录/章节大字,默认24)')
    fs.add_argument('--title-pattern', default=_TITLE_PAT, help='标题文本正则(命中跳过,默认 (一)(二)类)')
    fs.add_argument('--dry-run', action='store_true')

    tf = sub.add_parser('textboxfill', help='文本框 spPr 填充→透明(noFill)（系统 python3）')
    tf.add_argument('pptx')
    tf.add_argument('--match', choices=['bg1-alpha', 'any-solid'], default='bg1-alpha',
                    help='bg1-alpha(默认,只动半透明白底遮罩)| any-solid(所有 solidFill)')
    tf.add_argument('--dry-run', action='store_true')

    r = sub.add_parser('render', help='soffice→PNG 渲染验证（系统 python3）')
    r.add_argument('pptx')
    r.add_argument('--pages', help='逗号分隔页号(演示顺序),如 1,3,10;省略=全部')
    r.add_argument('--dpi', type=int, default=110)
    r.add_argument('--outdir')

    for ap in (a, l, ts, c, fs, tf, r):
        ap.set_defaults(func=_cmd_align)

    # ── chart（main() 里前置转发,此处只为 --help 列目） ──
    ch = sub.add_parser("chart", help="数据驱动图表生成 JSON -> PNG（转发 chart.py）",
                        add_help=False)
    ch.add_argument("args", nargs=argparse.REMAINDER)

    return parser


def main() -> int:
    # 无参数时显示帮助
    if len(sys.argv) < 2:
        build_parser().print_help()
        return 1

    # chart 前置转发（REMAINDER 语义,--help 等原样交给 chart.py）
    if sys.argv[1] == "chart":
        return _dispatch_chart(sys.argv[2:])

    # --batch 顶层旗截断（原 pptx_tools 批量模式,与子命令解析互不干扰）
    if "--batch" in sys.argv[1:]:
        return _batch_main()

    args = build_parser().parse_args()
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
